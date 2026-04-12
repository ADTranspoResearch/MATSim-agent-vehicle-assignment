from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_DIR = Path("/Users/natomanzolli/Documents/GitHub/MATSim-agent-vehicle-assignment/adoption prediction model")
OUTPUT_DIR = PROJECT_DIR / "validation_outputs" / "replacement_dynamics_presentation"
PROV_DIR = PROJECT_DIR / "validation_outputs" / "replacement_dynamics_ev_calibrated_model"
FSA_DIR = PROJECT_DIR / "validation_outputs" / "fsa_replacement_dynamics_ev_calibrated_model"
DATASET_DIR = PROJECT_DIR / "datasets"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_bullet_slide(prs, title, bullets, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
    if footer:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.0), Inches(0.4))
        p = tx.text_frame.paragraphs[0]
        p.text = footer
        p.font.size = Pt(11)
    return slide


def add_picture_slide(prs, title, image_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.1), width=Inches(12.0))
    if caption:
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
    return slide


def add_two_picture_slide(prs, title, image_left, image_right, caption_left="", caption_right=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_left), Inches(0.4), Inches(1.2), width=Inches(6.0))
    slide.shapes.add_picture(str(image_right), Inches(6.8), Inches(1.2), width=Inches(6.0))
    for left, top, text in [
        (0.5, 6.5, caption_left),
        (6.9, 6.5, caption_right),
    ]:
        tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(5.8), Inches(0.4))
        p = tx.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
    return slide


def add_table_slide(prs, title, dataframe, left=0.5, top=1.5, width=12.0, height=4.5, font_size=14):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows, cols = dataframe.shape[0] + 1, dataframe.shape[1] + 1
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    table.columns[0].width = Inches(2.0)
    for j, col in enumerate(dataframe.columns, start=1):
        table.cell(0, j).text = str(col)
    for i, idx in enumerate(dataframe.index, start=1):
        table.cell(i, 0).text = str(idx)
        for j, col in enumerate(dataframe.columns, start=1):
            val = dataframe.loc[idx, col]
            if isinstance(val, float):
                if abs(val) >= 1000:
                    text = f"{val:,.0f}"
                else:
                    text = f"{val:.3f}"
            else:
                text = str(val)
            table.cell(i, j).text = text
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                if r == 0:
                    p.font.bold = True
    return slide


def add_formula_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Model Math"
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.2), Inches(12.0), Inches(5.6))
    tf = box.text_frame
    tf.clear()
    lines = [
        "Province sales benchmark:",
        "s(k,t) = calibrated external sales share for vehicle type k in year t",
        "",
        "Province total fleet target:",
        "N(t) = N(t-1) * (1 + g(t)), where g(t) is the population-linked net growth rate",
        "",
        "Turnover and stock update:",
        "Entries(t) = tau(t) * N(t-1)",
        "Exits(t) = max(N(t-1) + Entries(t) - N(t), 0)",
        "Count(k,t) = Count(k,t-1) + Entries(t)*s(k,t) - Exits(t)*r(k,t)",
        "",
        "FSA shrinkage rule:",
        "w(FSA) = avg_entries / (avg_entries + K)",
        "sales_share(FSA,t) = w(FSA)*local(FSA,t) + (1-w(FSA))*province(t)",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20 if ":" in line and line else 18)
        if ":" in line and line:
            p.font.bold = True
    return slide


def build_stock_controls_chart(path):
    controls = pd.read_csv(PROV_DIR / "future_stock_controls.csv", index_col=0)
    fig, ax1 = plt.subplots(figsize=(10, 5))
    ax1.plot(controls.index, controls["target_total_fleet"], color="#1f4e79", linewidth=2.5)
    ax1.set_title("Population-Linked Province Fleet Target")
    ax1.set_xlabel("Year")
    ax1.set_ylabel("Target total fleet")
    ax2 = ax1.twinx()
    ax2.plot(controls.index, 100 * controls["future_population_growth"], color="#70ad47", linestyle="--", linewidth=2)
    ax2.set_ylabel("Population growth (%)")
    fig.tight_layout()
    fig.savefig(path, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return controls


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    province_summary = pd.read_csv(PROV_DIR / "population_linked_calibrated_2035_summary.csv", index_col=0)
    qc_summary = pd.read_csv(FSA_DIR / "quebec_city_summary.csv", index_col=0)
    fsa_summary = pd.read_csv(FSA_DIR / "fsa_population_linked_summary.csv")
    controls = build_stock_controls_chart(OUTPUT_DIR / "province_stock_controls.png")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Vehicle Adoption Forecasting",
        "Replacement-dynamics model with external market calibration, population-linked fleet growth, and FSA shrinkage",
    )

    add_bullet_slide(
        prs,
        "Datasets Used",
        [
            "SAAQ cached fleet, entry, and exit tables by year, vehicle type, and FSA",
            "External new vehicle registrations: Fig1-NMVRegist.xlsx (2017 to 2025 Q1)",
            "Quebec population reference: 1710000901-eng.csv (Q1 2015 to Q1 2026)",
            "Outputs from the calibrated province model and the calibrated FSA model",
        ],
        footer="All files are stored inside the adoption prediction model folder so the workflow is self-contained.",
    )

    add_bullet_slide(
        prs,
        "Core Assumptions",
        [
            "Total fleet size should evolve smoothly and follow population-linked growth, not unconstrained entry/exit extrapolation.",
            "New sales mix can change faster than the fleet mix because stock turnover is gradual.",
            "External registrations are used to calibrate the full future sales mix, not only EVs.",
            "At FSA level, sparse local data should borrow strength from the province benchmark through shrinkage.",
        ],
    )

    add_formula_slide(prs)

    add_bullet_slide(
        prs,
        "Why The Model Changed",
        [
            "Earlier versions allowed exits to stay above entries year after year, which made the total fleet collapse unrealistically.",
            "The updated model separates turnover from net fleet growth.",
            "Turnover comes from SAAQ entry intensity; net fleet growth comes from the Quebec population reference.",
            "This keeps province totals stable while still allowing the technology mix to evolve.",
        ],
    )

    add_picture_slide(
        prs,
        "Province Stock Controls",
        OUTPUT_DIR / "province_stock_controls.png",
        caption=(
            f"Future turnover rate is about {controls['future_turnover_rate'].iloc[-1]*100:.2f}% and "
            f"long-run population-linked growth is about {controls['future_population_growth'].iloc[-1]*100:.2f}% per year."
        ),
    )

    add_picture_slide(
        prs,
        "Province Results: Calibrated Sales Share",
        PROV_DIR / "population_linked_calibrated_sales_share.png",
        caption="Sales share changes faster than fleet share because only part of the stock is replaced each year.",
    )

    add_two_picture_slide(
        prs,
        "Province Results: Fleet Share and Counts",
        PROV_DIR / "population_linked_calibrated_fleet_market_share.png",
        PROV_DIR / "population_linked_calibrated_total_vehicle_counts.png",
        caption_left="Population-linked fleet market share by type",
        caption_right="Population-linked vehicle counts by type",
    )

    province_table = province_summary.copy()
    province_table["market_share_2035"] = 100 * province_table["market_share_2035"]
    province_table.columns = ["Count 2035", "Fleet Share 2035 (%)"]
    add_table_slide(prs, "Province 2035 Summary", province_table)

    add_bullet_slide(
        prs,
        "FSA-Level Extension",
        [
            "Each FSA keeps local fleet, entry, and exit history from the SAAQ cache.",
            "Future local sales share is a weighted combination of local trend and province benchmark.",
            "Shrinkage weight: w = avg_entries / (avg_entries + K), with K = 500 in the current notebook.",
            "Local total fleet is tied to a stable share of the province total instead of noisy local linear extrapolation.",
        ],
    )

    add_picture_slide(
        prs,
        "Selected FSAs: EV Fleet Share",
        FSA_DIR / "selected_fsa_ev_fleet_share.png",
        caption="Example comparison across selected postal areas using the shrinkage-based FSA model.",
    )

    add_two_picture_slide(
        prs,
        "Quebec City Aggregate Results",
        FSA_DIR / "quebec_city_sales_share.png",
        FSA_DIR / "quebec_city_fleet_market_share.png",
        caption_left="Quebec City sales share by type",
        caption_right="Quebec City fleet market share by type",
    )

    qc_table = qc_summary.copy()
    qc_table[["sales_share_2021", "sales_share_2035", "fleet_share_2021", "fleet_share_2035"]] *= 100
    qc_table.columns = [
        "Sales 2021 (%)",
        "Sales 2035 (%)",
        "Fleet 2021 (%)",
        "Fleet 2035 (%)",
        "Count 2035",
    ]
    add_table_slide(prs, "Quebec City 2035 Evaluation", qc_table, top=1.3, height=5.2, font_size=13)

    top_fsas = (
        fsa_summary.sort_values("electric_fleet_share_2035", ascending=False)[
            ["fsa", "avg_entries_per_year", "shrinkage_weight", "electric_fleet_share_2035"]
        ]
        .head(8)
        .copy()
    )
    top_fsas["electric_fleet_share_2035"] *= 100
    top_fsas = top_fsas.set_index("fsa")
    top_fsas.columns = ["Avg entries/year", "Shrinkage weight", "EV fleet share 2035 (%)"]
    add_table_slide(prs, "Highest 2035 EV Share FSAs", top_fsas, top=1.6, height=4.8, font_size=14)

    add_bullet_slide(
        prs,
        "Key Takeaways",
        [
            "The province model is now anchored to a realistic total-fleet path and a market-calibrated sales mix.",
            "The 2035 province EV fleet share is about 15.3%, with EV sales share around 19.0%.",
            "The FSA model is feasible with the current cache because hierarchical shrinkage prevents sparse areas from exploding.",
            "Quebec City remains SUV-heavy in both sales and fleet, while EV and HEV shares rise gradually through stock replacement.",
        ],
    )

    deck_path = OUTPUT_DIR / "replacement_dynamics_adoption_model_explainer.pptx"
    prs.save(deck_path)
    print(deck_path)


if __name__ == "__main__":
    main()
